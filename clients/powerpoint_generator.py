"""
PowerPoint Generator - Create Engaging, Animated Presentations
Builds .pptx files with advanced layouts and animations
"""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR, MSO_AUTO_SIZE
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from typing import List, Dict, Optional, Tuple
import re


# ============================================================================
# COLOR SCHEME (matching our brand)
# ============================================================================

COLORS = {
    'bg_primary': RGBColor(10, 25, 47),      # Navy
    'bg_secondary': RGBColor(17, 34, 64),    # Lighter navy
    'bg_card': RGBColor(26, 41, 64),         # Card background
    'accent_primary': RGBColor(100, 255, 218),  # Teal
    'accent_secondary': RGBColor(79, 209, 197), # Lighter teal
    'text_primary': RGBColor(204, 214, 246),    # Light gray
    'text_secondary': RGBColor(136, 146, 176),  # Medium gray
    'text_muted': RGBColor(73, 86, 112),        # Dark gray
    'white': RGBColor(255, 255, 255),
    'divider': RGBColor(35, 53, 84),
}


# ============================================================================
# HELPER FUNCTIONS
# ============================================================================

def extract_numbers(text: str) -> List[Tuple[str, str]]:
    """Extract numbers and their context from text"""
    pattern = r'(\d+(?:\.\d+)?)\s*([a-zA-Z%]+)?'
    matches = re.findall(pattern, text)

    results = []
    for num, unit in matches:
        if len(num) == 1 and not unit:
            continue
        results.append((num, unit or ''))

    return results[:6]


def detect_slide_type(title: str, bullets: List[str]) -> str:
    """Detect the best slide type based on content"""
    content = (title + ' ' + ' '.join(bullets)).lower()

    # Check for question
    if '?' in title or any('?' in b for b in bullets):
        return 'question_answer'

    # Check for timeline indicators
    if any(word in content for word in ['timeline', 'sequence', 'steps', 'phases', 'stage']):
        return 'journey'

    # Check for comparison indicators
    if any(word in content for word in ['vs', 'versus', 'before', 'after', 'comparison', 'compare']):
        return 'comparison'

    # Check for stats/numbers
    numbers = extract_numbers(content)
    if len(numbers) >= 2:
        # Check if this should be a hero stat (single dominant number)
        if len(bullets) <= 2 and len(numbers) == 1:
            return 'hero_stat'
        return 'stats'

    # Default to enhanced bullets
    return 'bullets'


def add_animation_fade_in(shape, delay=0):
    """Add fade-in entrance animation to a shape"""
    try:
        # Get the slide's timeline
        slide = shape.part.slide

        # Note: python-pptx doesn't have full animation API support
        # We'll add basic structure that PowerPoint will recognize
        # Full animations require XML manipulation

        # For now, we'll prepare the structure
        # Users can enhance in PowerPoint if needed
        pass
    except:
        pass


# ============================================================================
# POWERPOINT PRESENTATION CLASS
# ============================================================================

class EngagingPresentation:
    """Create engaging PowerPoint presentations with advanced layouts"""

    def __init__(self, title: str = "Presentation"):
        """Initialize a new presentation"""
        self.prs = Presentation()
        self.prs.slide_width = Inches(10)  # 16:9 aspect ratio
        self.prs.slide_height = Inches(5.625)
        self.title = title

        # Define slide layouts (we'll use blank and customize)
        self.blank_layout = self.prs.slide_layouts[6]  # Blank layout

    def add_title_slide(self, title: str, subtitle: str = ""):
        """Add an engaging title slide"""
        slide = self.prs.slides.add_slide(self.blank_layout)

        # Background
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = COLORS['bg_primary']

        # Accent bar at top
        accent_bar = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(0), Inches(0),
            Inches(10), Inches(0.05)
        )
        accent_bar.fill.solid()
        accent_bar.fill.fore_color.rgb = COLORS['accent_primary']
        accent_bar.line.fill.background()

        # Title
        title_box = slide.shapes.add_textbox(
            Inches(1), Inches(1.5),
            Inches(8), Inches(1.5)
        )
        title_frame = title_box.text_frame
        title_frame.text = title
        title_frame.word_wrap = True

        # Format title
        p = title_frame.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        p.font.size = Pt(54)
        p.font.bold = True
        p.font.color.rgb = COLORS['text_primary']

        # Subtitle if provided
        if subtitle:
            subtitle_box = slide.shapes.add_textbox(
                Inches(1), Inches(3.2),
                Inches(8), Inches(1)
            )
            subtitle_frame = subtitle_box.text_frame
            subtitle_frame.text = subtitle

            p = subtitle_frame.paragraphs[0]
            p.alignment = PP_ALIGN.CENTER
            p.font.size = Pt(24)
            p.font.color.rgb = COLORS['text_secondary']

        return slide

    def add_hero_stat_slide(self, big_number: str, unit: str, context: str, supporting: str = ""):
        """Add a hero stat slide with massive number"""
        slide = self.prs.slides.add_slide(self.blank_layout)

        # Background
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = COLORS['bg_primary']

        # Big number (massive!)
        number_box = slide.shapes.add_textbox(
            Inches(1), Inches(1),
            Inches(8), Inches(2)
        )
        number_frame = number_box.text_frame
        number_frame.text = big_number

        p = number_frame.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        p.font.size = Pt(120)
        p.font.bold = True
        p.font.color.rgb = COLORS['accent_primary']

        # Unit (if provided)
        if unit:
            unit_box = slide.shapes.add_textbox(
                Inches(1), Inches(3),
                Inches(8), Inches(0.5)
            )
            unit_frame = unit_box.text_frame
            unit_frame.text = unit

            p = unit_frame.paragraphs[0]
            p.alignment = PP_ALIGN.CENTER
            p.font.size = Pt(36)
            p.font.color.rgb = COLORS['text_secondary']

        # Context
        context_box = slide.shapes.add_textbox(
            Inches(1.5), Inches(3.8),
            Inches(7), Inches(0.8)
        )
        context_frame = context_box.text_frame
        context_frame.text = context
        context_frame.word_wrap = True

        p = context_frame.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        p.font.size = Pt(28)
        p.font.color.rgb = COLORS['text_primary']

        # Supporting text
        if supporting:
            support_box = slide.shapes.add_textbox(
                Inches(2), Inches(4.8),
                Inches(6), Inches(0.5)
            )
            support_frame = support_box.text_frame
            support_frame.text = supporting

            p = support_frame.paragraphs[0]
            p.alignment = PP_ALIGN.CENTER
            p.font.size = Pt(18)
            p.font.color.rgb = COLORS['text_secondary']

        return slide

    def add_journey_slide(self, title: str, steps: List[str]):
        """Add a journey/timeline slide showing progression"""
        slide = self.prs.slides.add_slide(self.blank_layout)

        # Background
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = COLORS['bg_primary']

        # Title
        self._add_slide_title(slide, title)

        # Journey visualization
        num_steps = min(len(steps), 5)
        step_width = 7 / num_steps
        start_x = 1.5
        y_pos = 2.2

        for i, step in enumerate(steps[:5]):
            x_pos = start_x + (i * step_width)

            # Circle for step
            circle = slide.shapes.add_shape(
                MSO_SHAPE.OVAL,
                Inches(x_pos + step_width/2 - 0.2), Inches(y_pos),
                Inches(0.4), Inches(0.4)
            )
            circle.fill.solid()
            circle.fill.fore_color.rgb = COLORS['accent_primary']
            circle.line.fill.background()

            # Arrow to next (if not last)
            if i < num_steps - 1:
                arrow = slide.shapes.add_shape(
                    MSO_SHAPE.RIGHT_ARROW,
                    Inches(x_pos + step_width/2 + 0.25), Inches(y_pos + 0.15),
                    Inches(step_width - 0.5), Inches(0.1)
                )
                arrow.fill.solid()
                arrow.fill.fore_color.rgb = COLORS['accent_secondary']
                arrow.line.fill.background()

            # Step text
            text_box = slide.shapes.add_textbox(
                Inches(x_pos), Inches(y_pos + 0.6),
                Inches(step_width), Inches(1)
            )
            text_frame = text_box.text_frame
            text_frame.text = step
            text_frame.word_wrap = True

            p = text_frame.paragraphs[0]
            p.alignment = PP_ALIGN.CENTER
            p.font.size = Pt(14)
            p.font.color.rgb = COLORS['text_primary']

        return slide

    def add_comparison_slide(self, title: str, left_items: List[str], right_items: List[str],
                           left_label: str = "Before", right_label: str = "After"):
        """Add a before/after comparison slide"""
        slide = self.prs.slides.add_slide(self.blank_layout)

        # Background
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = COLORS['bg_primary']

        # Title
        self._add_slide_title(slide, title)

        # Divider line
        divider = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(5), Inches(1.2),
            Inches(0.02), Inches(3.8)
        )
        divider.fill.solid()
        divider.fill.fore_color.rgb = COLORS['divider']
        divider.line.fill.background()

        # Left column header
        left_header = slide.shapes.add_textbox(
            Inches(0.5), Inches(1.3),
            Inches(4), Inches(0.4)
        )
        left_frame = left_header.text_frame
        left_frame.text = left_label

        p = left_frame.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        p.font.size = Pt(24)
        p.font.bold = True
        p.font.color.rgb = COLORS['accent_primary']

        # Right column header
        right_header = slide.shapes.add_textbox(
            Inches(5.5), Inches(1.3),
            Inches(4), Inches(0.4)
        )
        right_frame = right_header.text_frame
        right_frame.text = right_label

        p = right_frame.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        p.font.size = Pt(24)
        p.font.bold = True
        p.font.color.rgb = COLORS['accent_primary']

        # Left items
        self._add_bullet_column(slide, left_items, Inches(0.5), Inches(2), Inches(4))

        # Right items
        self._add_bullet_column(slide, right_items, Inches(5.5), Inches(2), Inches(4))

        return slide

    def add_stats_slide(self, title: str, stats: List[Tuple[str, str, str]]):
        """
        Add a stats slide with card layout
        stats: List of (value, unit, label) tuples
        """
        slide = self.prs.slides.add_slide(self.blank_layout)

        # Background
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = COLORS['bg_primary']

        # Title
        self._add_slide_title(slide, title)

        # Stat cards
        num_stats = min(len(stats), 4)
        card_width = 1.8
        spacing = 0.3
        total_width = (num_stats * card_width) + ((num_stats - 1) * spacing)
        start_x = (10 - total_width) / 2

        for i, (value, unit, label) in enumerate(stats[:4]):
            x_pos = start_x + (i * (card_width + spacing))

            # Card background
            card = slide.shapes.add_shape(
                MSO_SHAPE.ROUNDED_RECTANGLE,
                Inches(x_pos), Inches(2),
                Inches(card_width), Inches(2.5)
            )
            card.fill.solid()
            card.fill.fore_color.rgb = COLORS['bg_card']
            card.line.color.rgb = COLORS['accent_primary']
            card.line.width = Pt(2)

            # Value
            value_box = slide.shapes.add_textbox(
                Inches(x_pos + 0.1), Inches(2.4),
                Inches(card_width - 0.2), Inches(1)
            )
            value_frame = value_box.text_frame
            value_frame.text = value + unit

            p = value_frame.paragraphs[0]
            p.alignment = PP_ALIGN.CENTER
            p.font.size = Pt(48)
            p.font.bold = True
            p.font.color.rgb = COLORS['accent_primary']

            # Label
            label_box = slide.shapes.add_textbox(
                Inches(x_pos + 0.1), Inches(3.6),
                Inches(card_width - 0.2), Inches(0.7)
            )
            label_frame = label_box.text_frame
            label_frame.text = label
            label_frame.word_wrap = True

            p = label_frame.paragraphs[0]
            p.alignment = PP_ALIGN.CENTER
            p.font.size = Pt(14)
            p.font.color.rgb = COLORS['text_secondary']

        return slide

    def add_bullets_slide(self, title: str, bullets: List[str]):
        """Add an enhanced bullet point slide"""
        slide = self.prs.slides.add_slide(self.blank_layout)

        # Background
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = COLORS['bg_primary']

        # Title
        self._add_slide_title(slide, title)

        # Bullets with numbering
        y_start = 2
        bullet_spacing = 0.7

        for i, bullet in enumerate(bullets[:6]):
            y_pos = y_start + (i * bullet_spacing)

            # Number circle
            circle = slide.shapes.add_shape(
                MSO_SHAPE.OVAL,
                Inches(1), Inches(y_pos),
                Inches(0.4), Inches(0.4)
            )
            circle.fill.solid()
            circle.fill.fore_color.rgb = COLORS['accent_primary']
            circle.fill.fore_color.brightness = -0.3
            circle.line.fill.background()

            # Number
            num_box = slide.shapes.add_textbox(
                Inches(1), Inches(y_pos),
                Inches(0.4), Inches(0.4)
            )
            num_frame = num_box.text_frame
            num_frame.text = str(i + 1)
            num_frame.vertical_anchor = MSO_ANCHOR.MIDDLE

            p = num_frame.paragraphs[0]
            p.alignment = PP_ALIGN.CENTER
            p.font.size = Pt(16)
            p.font.bold = True
            p.font.color.rgb = COLORS['accent_primary']

            # Bullet text
            text_box = slide.shapes.add_textbox(
                Inches(1.6), Inches(y_pos - 0.05),
                Inches(7.5), Inches(0.6)
            )
            text_frame = text_box.text_frame
            text_frame.text = bullet
            text_frame.word_wrap = True

            p = text_frame.paragraphs[0]
            p.font.size = Pt(18)
            p.font.color.rgb = COLORS['text_primary']

        return slide

    def add_question_answer_slide(self, question: str, answer: str, context: str = ""):
        """Add a question-answer slide that builds curiosity"""
        slide = self.prs.slides.add_slide(self.blank_layout)

        # Background
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = COLORS['bg_primary']

        # Question
        question_box = slide.shapes.add_textbox(
            Inches(1.5), Inches(1),
            Inches(7), Inches(1.2)
        )
        question_frame = question_box.text_frame
        question_frame.text = question
        question_frame.word_wrap = True

        p = question_frame.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        p.font.size = Pt(32)
        p.font.bold = True
        p.font.color.rgb = COLORS['text_secondary']

        # Answer box (highlighted)
        answer_bg = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(2), Inches(2.5),
            Inches(6), Inches(1.2)
        )
        answer_bg.fill.solid()
        answer_bg.fill.fore_color.rgb = COLORS['bg_card']
        answer_bg.line.color.rgb = COLORS['accent_primary']
        answer_bg.line.width = Pt(3)

        # Answer text
        answer_box = slide.shapes.add_textbox(
            Inches(2.2), Inches(2.7),
            Inches(5.6), Inches(0.8)
        )
        answer_frame = answer_box.text_frame
        answer_frame.text = answer
        answer_frame.word_wrap = True

        p = answer_frame.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        p.font.size = Pt(48)
        p.font.bold = True
        p.font.color.rgb = COLORS['accent_primary']

        # Context
        if context:
            context_box = slide.shapes.add_textbox(
                Inches(2), Inches(4),
                Inches(6), Inches(0.8)
            )
            context_frame = context_box.text_frame
            context_frame.text = context
            context_frame.word_wrap = True

            p = context_frame.paragraphs[0]
            p.alignment = PP_ALIGN.CENTER
            p.font.size = Pt(20)
            p.font.color.rgb = COLORS['text_secondary']

        return slide

    def _add_slide_title(self, slide, title: str):
        """Helper: Add title to a slide"""
        title_box = slide.shapes.add_textbox(
            Inches(0.5), Inches(0.3),
            Inches(9), Inches(0.6)
        )
        title_frame = title_box.text_frame
        title_frame.text = title

        p = title_frame.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        p.font.size = Pt(36)
        p.font.bold = True
        p.font.color.rgb = COLORS['text_primary']

        # Underline
        underline = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(3.5), Inches(1),
            Inches(3), Inches(0.03)
        )
        underline.fill.solid()
        underline.fill.fore_color.rgb = COLORS['accent_primary']
        underline.line.fill.background()

    def _add_bullet_column(self, slide, items: List[str], left: Inches, top: Inches, width: Inches):
        """Helper: Add a column of bullet points"""
        spacing = 0.5
        for i, item in enumerate(items[:5]):
            y_pos = top + (i * spacing)

            text_box = slide.shapes.add_textbox(
                left, Inches(y_pos),
                width, Inches(0.4)
            )
            text_frame = text_box.text_frame
            text_frame.text = "• " + item
            text_frame.word_wrap = True

            p = text_frame.paragraphs[0]
            p.font.size = Pt(16)
            p.font.color.rgb = COLORS['text_primary']

    def save(self, filename: str):
        """Save the presentation"""
        if not filename.endswith('.pptx'):
            filename += '.pptx'
        self.prs.save(filename)
        return filename

    def add_image_slide(self, image_path: str):
        """Add a full-bleed image slide."""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        slide.shapes.add_picture(
            image_path,
            Inches(0),
            Inches(0),
            width=self.prs.slide_width,
            height=self.prs.slide_height
        )


# ============================================================================
# CONVENIENCE FUNCTIONS
# ============================================================================

def create_presentation_from_slides(slides_data: List[Dict], title: str = "Presentation") -> str:
    """
    Create a PowerPoint presentation from slide data

    slides_data: List of dicts with keys:
        - title: str
        - bullets: List[str]
        - type: Optional[str] - slide type override

    Returns: filename of saved presentation
    """
    prs = EngagingPresentation(title)

    # Add title slide
    prs.add_title_slide(title, "Generated with QLSV2 Learning Studio")

    # Add content slides
    for slide_data in slides_data:
        image_path = slide_data.get('image_path')
        if image_path:
            prs.add_image_slide(image_path)
            continue
        slide_title = slide_data.get('title', 'Slide')
        bullets = slide_data.get('bullets', [])
        slide_type = slide_data.get('type') or detect_slide_type(slide_title, bullets)

        if slide_type == 'hero_stat':
            # Extract number for hero stat
            numbers = extract_numbers(slide_title + ' ' + ' '.join(bullets))
            if numbers:
                value, unit = numbers[0]
                context = slide_title
                supporting = bullets[0] if bullets else ""
                prs.add_hero_stat_slide(value, unit, context, supporting)
            else:
                # Fallback to bullets
                prs.add_bullets_slide(slide_title, bullets)

        elif slide_type == 'journey':
            prs.add_journey_slide(slide_title, bullets)

        elif slide_type == 'comparison':
            mid_point = len(bullets) // 2
            left = bullets[:mid_point]
            right = bullets[mid_point:]
            prs.add_comparison_slide(slide_title, left, right)

        elif slide_type == 'stats':
            # Extract stats
            stats = []
            for bullet in bullets[:4]:
                numbers = extract_numbers(bullet)
                if numbers:
                    value, unit = numbers[0]
                    label = bullet.replace(value, '').replace(unit, '').strip()
                    stats.append((value, unit, label[:30]))

            if stats:
                prs.add_stats_slide(slide_title, stats)
            else:
                prs.add_bullets_slide(slide_title, bullets)

        elif slide_type == 'question_answer':
            question = slide_title if '?' in slide_title else bullets[0] if bullets else "Question?"
            answer = bullets[1] if len(bullets) > 1 else bullets[0] if bullets else "Answer"
            context = bullets[2] if len(bullets) > 2 else ""
            prs.add_question_answer_slide(question, answer, context)

        else:  # bullets
            prs.add_bullets_slide(slide_title, bullets)

    # Save
    filename = f"presentation_{title.replace(' ', '_')[:30]}.pptx"
    return prs.save(filename)
